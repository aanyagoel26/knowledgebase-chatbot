from fastapi import FastAPI, Request, BackgroundTasks, Response, HTTPException
from fastapi.responses import FileResponse
from pydantic import BaseModel
import psycopg2
from psycopg2.extras import execute_values
import requests
import os
import shutil
import hashlib
import re
import csv
import threading
import time
import uuid
from datetime import datetime
import fitz
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation


app = FastAPI()


DB_HOST = "localhost"
DB_NAME = "kb_chatbot"
DB_USER = "postgres"
DB_PASSWORD = "Aanya2612"

CHAT_MODEL = "qwen2.5:7b"
EMBEDDING_MODEL = "nomic-embed-text"

DEFAULT_KNOWLEDGE_BASE_FOLDER = "knowledge_base"
KNOWLEDGE_BASE_FOLDER = DEFAULT_KNOWLEDGE_BASE_FOLDER
UPLOAD_FOLDER = KNOWLEDGE_BASE_FOLDER
UI_FILE = "kb_chat.html"

AUTO_SCAN_INTERVAL_SECONDS = 3600

MAX_CHUNK_SIZE = 1500
CHUNK_OVERLAP = 200

PER_DOCUMENT_VECTOR_TOP_K = 6
PER_DOCUMENT_KEYWORD_TOP_K = 6
PER_DOCUMENT_DIRECT_TOP_K = 2
PER_DOCUMENT_CONTEXT_LIMIT = 5

NEIGHBOR_WINDOW = 1
MAX_CONTEXT_CHUNKS_TOTAL = 14

EMBEDDING_BATCH_SIZE = 32

scan_lock = threading.Lock()


class RetrieveRequest(BaseModel):
    question: str
    document_ids: list[int] = []


class ChatRequest(BaseModel):
    question: str
    session_id: int | None = None
    document_ids: list[int] = []


class LoginRequest(BaseModel):
    email: str
    password: str


class FolderRequest(BaseModel):
    folder_path: str


def get_db_connection():
    return psycopg2.connect(
        host=DB_HOST,
        database=DB_NAME,
        user=DB_USER,
        password=DB_PASSWORD
    )


def hash_password(password):
    return hashlib.sha256(password.encode("utf-8")).hexdigest()


def get_session_token(request: Request):
    return request.cookies.get("kb_session_token")


def require_login(request: Request):
    token = get_session_token(request)

    if not token:
        raise HTTPException(status_code=401, detail="Not logged in")

    conn = get_db_connection()
    cursor = conn.cursor()

    cursor.execute(
        """
        SELECT e.employee_id, e.name, e.email, e.role
        FROM employee_sessions s
        JOIN employees e
        ON s.employee_id = e.employee_id
        WHERE s.session_token = %s
        """,
        (token,)
    )

    row = cursor.fetchone()

    cursor.close()
    conn.close()

    if not row:
        raise HTTPException(status_code=401, detail="Invalid session")

    return {
        "employee_id": row[0],
        "name": row[1],
        "email": row[2],
        "role": row[3]
    }


def ensure_schema_updates():
    conn = get_db_connection()
    cursor = conn.cursor()

    cursor.execute(
        """
        ALTER TABLE documents
        ADD COLUMN IF NOT EXISTS indexing_status TEXT DEFAULT 'ready';
        """
    )

    cursor.execute(
        """
        ALTER TABLE documents
        ADD COLUMN IF NOT EXISTS error_message TEXT;
        """
    )

    cursor.execute(
        """
        ALTER TABLE documents
        ADD COLUMN IF NOT EXISTS chunk_count INTEGER DEFAULT 0;
        """
    )

    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS employees (
            employee_id SERIAL PRIMARY KEY,
            name TEXT NOT NULL,
            email TEXT UNIQUE NOT NULL,
            password_hash TEXT NOT NULL,
            role TEXT DEFAULT 'employee',
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        );
        """
    )

    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS employee_sessions (
            session_id SERIAL PRIMARY KEY,
            employee_id INTEGER REFERENCES employees(employee_id) ON DELETE CASCADE,
            session_token TEXT UNIQUE NOT NULL,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        );
        """
    )

    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS app_settings (
            setting_key TEXT PRIMARY KEY,
            setting_value TEXT NOT NULL,
            updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        );
        """
    )

    cursor.execute(
        """
        SELECT employee_id
        FROM employees
        WHERE email = %s
        """,
        ("admin@motherson.com",)
    )

    existing_admin = cursor.fetchone()

    if not existing_admin:
        cursor.execute(
            """
            INSERT INTO employees(name, email, password_hash, role)
            VALUES (%s, %s, %s, %s)
            """,
            (
                "Admin",
                "admin@motherson.com",
                hash_password("admin123"),
                "admin"
            )
        )

    conn.commit()
    cursor.close()
    conn.close()


def get_setting(key, default_value=None):
    conn = get_db_connection()
    cursor = conn.cursor()

    cursor.execute(
        """
        SELECT setting_value
        FROM app_settings
        WHERE setting_key = %s
        """,
        (key,)
    )

    row = cursor.fetchone()

    cursor.close()
    conn.close()

    if row:
        return row[0]

    return default_value


def set_setting(key, value):
    conn = get_db_connection()
    cursor = conn.cursor()

    cursor.execute(
        """
        INSERT INTO app_settings(setting_key, setting_value)
        VALUES (%s, %s)
        ON CONFLICT(setting_key)
        DO UPDATE SET
            setting_value = EXCLUDED.setting_value,
            updated_at = CURRENT_TIMESTAMP
        """,
        (key, value)
    )

    conn.commit()
    cursor.close()
    conn.close()


def load_knowledge_folder_from_db():
    global KNOWLEDGE_BASE_FOLDER
    global UPLOAD_FOLDER

    saved_folder = get_setting("knowledge_base_folder", DEFAULT_KNOWLEDGE_BASE_FOLDER)

    KNOWLEDGE_BASE_FOLDER = saved_folder
    UPLOAD_FOLDER = KNOWLEDGE_BASE_FOLDER


def ensure_folders():
    os.makedirs(KNOWLEDGE_BASE_FOLDER, exist_ok=True)


def calculate_file_hash(file_path):
    sha256 = hashlib.sha256()

    with open(file_path, "rb") as file:
        for block in iter(lambda: file.read(1024 * 1024), b""):
            sha256.update(block)

    return sha256.hexdigest()


def get_file_metadata(file_path):
    stat = os.stat(file_path)

    return {
        "file_size": stat.st_size,
        "last_modified": datetime.fromtimestamp(stat.st_mtime)
    }


def clean_text(text):
    text = text.replace("\r", "")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def tokenize(text):
    text = text.lower()
    text = re.sub(r"[^a-z0-9\s]", " ", text)

    stopwords = {
        "the", "is", "are", "was", "were", "a", "an", "of", "for", "to",
        "in", "on", "and", "or", "with", "from", "by", "as", "at", "this",
        "that", "it", "be", "can", "what", "which", "who", "when", "where",
        "how", "tell", "me", "about", "give", "list", "explain", "show",
        "please", "kindly", "all", "any"
    }

    return [word for word in text.split() if word not in stopwords]


def extract_pdf(file_path):
    document = fitz.open(file_path)
    text = ""

    for page_number, page in enumerate(document, start=1):
        text += f"\n\nPage {page_number}\n{page.get_text()}"

    document.close()
    return text


def extract_docx(file_path):
    doc = Document(file_path)
    text = ""

    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            text += paragraph.text.strip() + "\n"

    for table in doc.tables:
        for row in table.rows:
            row_values = []

            for cell in row.cells:
                if cell.text.strip():
                    row_values.append(cell.text.strip())

            if row_values:
                text += " | ".join(row_values) + "\n"

    return text


def extract_xlsx(file_path):
    workbook = load_workbook(file_path, data_only=True)
    text = ""

    for sheet in workbook.worksheets:
        text += f"\n\nSheet: {sheet.title}\n"

        for row in sheet.iter_rows(values_only=True):
            values = []

            for cell in row:
                if cell is not None:
                    values.append(str(cell))

            if values:
                text += " | ".join(values) + "\n"

    return text


def extract_csv(file_path):
    text = ""

    with open(file_path, "r", encoding="utf-8", errors="ignore") as file:
        reader = csv.reader(file)

        for row in reader:
            text += " | ".join(row) + "\n"

    return text


def extract_pptx(file_path):
    presentation = Presentation(file_path)
    text = ""

    for slide_number, slide in enumerate(presentation.slides, start=1):
        text += f"\n\nSlide {slide_number}\n"

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text += shape.text.strip() + "\n"

    return text


def extract_text_file(file_path):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as file:
        return file.read()


def extract_text(file_path):
    extension = file_path.lower().split(".")[-1]

    if extension == "pdf":
        return extract_pdf(file_path)

    if extension == "docx":
        return extract_docx(file_path)

    if extension == "xlsx":
        return extract_xlsx(file_path)

    if extension == "csv":
        return extract_csv(file_path)

    if extension == "pptx":
        return extract_pptx(file_path)

    if extension in ["txt", "md"]:
        return extract_text_file(file_path)

    raise ValueError("Unsupported file type: " + extension)


def is_supported_file(file_path):
    supported_extensions = ["pdf", "docx", "xlsx", "csv", "pptx", "txt", "md"]
    extension = file_path.lower().split(".")[-1]
    return extension in supported_extensions


def recursive_split(text, separators):
    if len(text) <= MAX_CHUNK_SIZE:
        return [text]

    if not separators:
        return [
            text[i:i + MAX_CHUNK_SIZE]
            for i in range(0, len(text), MAX_CHUNK_SIZE)
        ]

    separator = separators[0]
    parts = text.split(separator)

    if len(parts) == 1:
        return recursive_split(text, separators[1:])

    chunks = []
    current = ""

    for part in parts:
        part = part.strip()

        if not part:
            continue

        candidate = current + separator + part if current else part

        if len(candidate) <= MAX_CHUNK_SIZE:
            current = candidate
        else:
            if current:
                chunks.extend(recursive_split(current, separators[1:]))

            current = part

    if current:
        chunks.extend(recursive_split(current, separators[1:]))

    return chunks


def add_overlap(chunks):
    final_chunks = []

    for chunk in chunks:
        chunk = chunk.strip()

        if not chunk:
            continue

        if not final_chunks:
            final_chunks.append(chunk)
        else:
            previous = final_chunks[-1]
            overlap_text = previous[-CHUNK_OVERLAP:]
            combined = overlap_text + "\n" + chunk

            if len(combined) <= MAX_CHUNK_SIZE + CHUNK_OVERLAP:
                final_chunks.append(combined.strip())
            else:
                final_chunks.append(chunk)

    return final_chunks


def split_text_into_chunks(text):
    text = clean_text(text)

    separators = [
        "\n\n",
        "\n",
        ". ",
        "; ",
        ", ",
        " "
    ]

    raw_chunks = recursive_split(text, separators)

    safe_chunks = []

    for chunk in raw_chunks:
        if len(chunk) <= MAX_CHUNK_SIZE:
            safe_chunks.append(chunk)
        else:
            for i in range(0, len(chunk), MAX_CHUNK_SIZE):
                safe_chunks.append(chunk[i:i + MAX_CHUNK_SIZE])

    final_chunks = add_overlap(safe_chunks)

    return [
        chunk.strip()
        for chunk in final_chunks
        if len(chunk.strip()) > 30
    ]


def generate_embedding(text):
    response = requests.post(
        "http://localhost:11434/api/embeddings",
        json={
            "model": EMBEDDING_MODEL,
            "prompt": text
        },
        timeout=120
    )

    if response.status_code != 200:
        print("Embedding failed")
        print(response.text)
        raise Exception("Embedding API failed")

    return response.json()["embedding"]


def generate_embeddings_batch(texts):
    if not texts:
        return []

    try:
        response = requests.post(
            "http://localhost:11434/api/embed",
            json={
                "model": EMBEDDING_MODEL,
                "input": texts
            },
            timeout=300
        )

        if response.status_code == 200:
            data = response.json()

            if "embeddings" in data:
                return data["embeddings"]

    except Exception as error:
        print("Batch embedding failed, falling back to single embeddings.")
        print(error)

    embeddings = []

    for text in texts:
        embeddings.append(generate_embedding(text))

    return embeddings


def update_document_status(document_id, status, error_message=None, chunk_count=None):
    conn = get_db_connection()
    cursor = conn.cursor()

    if status == "ready":
        cursor.execute(
            """
            UPDATE documents
            SET indexing_status = %s,
                error_message = %s,
                chunk_count = %s,
                indexed_at = CURRENT_TIMESTAMP,
                updated_at = CURRENT_TIMESTAMP
            WHERE document_id = %s
            """,
            (
                status,
                error_message,
                chunk_count,
                document_id
            )
        )
    else:
        cursor.execute(
            """
            UPDATE documents
            SET indexing_status = %s,
                error_message = %s,
                chunk_count = COALESCE(%s, chunk_count),
                updated_at = CURRENT_TIMESTAMP
            WHERE document_id = %s
            """,
            (
                status,
                error_message,
                chunk_count,
                document_id
            )
        )

    conn.commit()
    cursor.close()
    conn.close()


def queue_file_for_indexing(file_path, source_type, force_reindex=False):
    absolute_path = os.path.abspath(file_path)
    original_filename = os.path.basename(file_path)
    metadata = get_file_metadata(file_path)
    file_hash = calculate_file_hash(file_path)

    conn = get_db_connection()
    cursor = conn.cursor()

    cursor.execute(
        """
        SELECT document_id, file_hash, file_size, version, indexing_status
        FROM documents
        WHERE file_path = %s
        """,
        (absolute_path,)
    )

    existing_by_path = cursor.fetchone()

    if existing_by_path:
        document_id, old_hash, old_size, old_version, old_status = existing_by_path

        unchanged = (
            old_hash == file_hash
            and old_size == metadata["file_size"]
        )

        if unchanged and old_status == "ready" and not force_reindex:
            cursor.close()
            conn.close()

            return {
                "filename": original_filename,
                "document_id": document_id,
                "status": "skipped_unchanged",
                "message": "File already indexed and unchanged.",
                "scheduled": False
            }

        if old_status in ["pending", "indexing"] and not force_reindex:
            cursor.close()
            conn.close()

            return {
                "filename": original_filename,
                "document_id": document_id,
                "status": "already_processing",
                "message": "File is already being indexed.",
                "scheduled": False
            }

        new_version = old_version

        if old_hash != file_hash or force_reindex:
            new_version = old_version + 1

        cursor.execute(
            """
            UPDATE documents
            SET file_hash = %s,
                file_size = %s,
                last_modified = %s,
                version = %s,
                source_type = %s,
                indexing_status = 'pending',
                error_message = NULL,
                chunk_count = 0,
                updated_at = CURRENT_TIMESTAMP
            WHERE document_id = %s
            """,
            (
                file_hash,
                metadata["file_size"],
                metadata["last_modified"],
                new_version,
                source_type,
                document_id
            )
        )

        conn.commit()
        cursor.close()
        conn.close()

        return {
            "filename": original_filename,
            "document_id": document_id,
            "status": "queued_for_reindexing",
            "message": "Updated file detected. Re-indexing started.",
            "scheduled": True
        }

    cursor.execute(
        """
        SELECT document_id
        FROM documents
        WHERE file_hash = %s
        """,
        (file_hash,)
    )

    existing_by_hash = cursor.fetchone()

    if existing_by_hash and not force_reindex:
        cursor.close()
        conn.close()

        return {
            "filename": original_filename,
            "document_id": existing_by_hash[0],
            "status": "skipped_duplicate",
            "message": "This file already exists in the knowledge base.",
            "scheduled": False
        }

    cursor.execute(
        """
        INSERT INTO documents
        (
            original_filename,
            file_path,
            file_hash,
            file_size,
            last_modified,
            source_type,
            indexing_status,
            error_message,
            chunk_count
        )
        VALUES (%s, %s, %s, %s, %s, %s, 'pending', NULL, 0)
        RETURNING document_id
        """,
        (
            original_filename,
            absolute_path,
            file_hash,
            metadata["file_size"],
            metadata["last_modified"],
            source_type
        )
    )

    document_id = cursor.fetchone()[0]

    conn.commit()
    cursor.close()
    conn.close()

    return {
        "filename": original_filename,
        "document_id": document_id,
        "status": "queued_new_document",
        "message": "New file queued for background indexing.",
        "scheduled": True
    }


def process_document_indexing(document_id, file_path):
    print("\n" + "=" * 80)
    print("BACKGROUND INDEXING STARTED")
    print("=" * 80)
    print("Document ID:", document_id)
    print("File:", file_path)

    try:
        update_document_status(document_id, "indexing", None, 0)

        extracted_text = extract_text(file_path)
        chunks = split_text_into_chunks(extracted_text)

        conn = get_db_connection()
        cursor = conn.cursor()

        cursor.execute(
            """
            DELETE FROM document_chunks
            WHERE document_id = %s
            """,
            (document_id,)
        )

        inserted_count = 0

        for batch_start in range(0, len(chunks), EMBEDDING_BATCH_SIZE):
            batch_chunks = chunks[batch_start:batch_start + EMBEDDING_BATCH_SIZE]

            embeddings = generate_embeddings_batch(batch_chunks)

            rows = []

            for index, chunk in enumerate(batch_chunks):
                chunk_number = batch_start + index + 1
                embedding = embeddings[index]

                rows.append(
                    (
                        document_id,
                        chunk_number,
                        chunk,
                        str(embedding),
                        len(tokenize(chunk))
                    )
                )

            execute_values(
                cursor,
                """
                INSERT INTO document_chunks
                (document_id, chunk_number, content, embedding, token_count)
                VALUES %s
                """,
                rows,
                template="(%s, %s, %s, %s::vector, %s)"
            )

            conn.commit()
            inserted_count += len(rows)

        cursor.close()
        conn.close()

        update_document_status(document_id, "ready", None, inserted_count)

        print("BACKGROUND INDEXING COMPLETED")
        print("Document ID:", document_id)
        print("Chunks inserted:", inserted_count)

    except Exception as error:
        error_text = str(error)

        print("BACKGROUND INDEXING FAILED")
        print("Document ID:", document_id)
        print(error_text)

        update_document_status(document_id, "failed", error_text, 0)


def scan_knowledge_base_once(force_reindex=False):
    if scan_lock.locked():
        print("Scan already running. Skipping this cycle.")
        return []

    results = []

    with scan_lock:
        ensure_folders()

        print("\n" + "=" * 80)
        print("SCANNING KNOWLEDGE BASE FOLDER")
        print("=" * 80)
        print("Folder:", os.path.abspath(KNOWLEDGE_BASE_FOLDER))

        for root, dirs, files in os.walk(KNOWLEDGE_BASE_FOLDER):
            for filename in files:
                file_path = os.path.join(root, filename)

                if not is_supported_file(file_path):
                    continue

                result = queue_file_for_indexing(
                    file_path=file_path,
                    source_type="knowledge_base",
                    force_reindex=force_reindex
                )

                results.append(result)
                print(result)

                if result.get("scheduled"):
                    process_document_indexing(
                        result["document_id"],
                        os.path.abspath(file_path)
                    )

    return results


def hourly_knowledge_base_watcher():
    print("\nHourly knowledge base watcher started.")
    print("Watching folder:", os.path.abspath(KNOWLEDGE_BASE_FOLDER))

    while True:
        time.sleep(AUTO_SCAN_INTERVAL_SECONDS)

        try:
            scan_knowledge_base_once(False)
        except Exception as error:
            print("Hourly watcher error:")
            print(error)


def get_ready_document_ids():
    conn = get_db_connection()
    cursor = conn.cursor()

    cursor.execute(
        """
        SELECT document_id
        FROM documents
        WHERE indexing_status = 'ready'
        ORDER BY updated_at DESC
        """
    )

    rows = cursor.fetchall()

    cursor.close()
    conn.close()

    return [row[0] for row in rows]


def build_document_filter(document_ids):
    if not document_ids:
        return "", []

    placeholders = ",".join(["%s"] * len(document_ids))
    return f" AND c.document_id IN ({placeholders}) ", document_ids


def get_vector_rows(cursor, question, document_ids=None, limit=PER_DOCUMENT_VECTOR_TOP_K):
    question_embedding = generate_embedding(question)
    document_ids = document_ids or []

    filter_sql, filter_params = build_document_filter(document_ids)

    query = f"""
        SELECT 
            c.chunk_id,
            c.document_id,
            d.original_filename,
            c.chunk_number,
            c.content,
            c.embedding <=> %s::vector AS distance
        FROM document_chunks c
        JOIN documents d
        ON c.document_id = d.document_id
        WHERE d.indexing_status = 'ready'
        {filter_sql}
        ORDER BY c.embedding <=> %s::vector
        LIMIT %s
    """

    params = (
        [str(question_embedding)]
        + filter_params
        + [str(question_embedding), limit]
    )

    cursor.execute(query, params)

    return cursor.fetchall()


def get_keyword_rows(cursor, question, document_ids=None, limit=PER_DOCUMENT_KEYWORD_TOP_K):
    keywords = tokenize(question)
    document_ids = document_ids or []

    if not keywords:
        return []

    keyword_conditions = []
    params = []

    for keyword in keywords:
        keyword_conditions.append("c.content ILIKE %s")
        params.append("%" + keyword + "%")

    document_filter_sql, document_filter_params = build_document_filter(document_ids)

    query = f"""
        SELECT 
            c.chunk_id,
            c.document_id,
            d.original_filename,
            c.chunk_number,
            c.content,
            1.0 AS distance
        FROM document_chunks c
        JOIN documents d
        ON c.document_id = d.document_id
        WHERE d.indexing_status = 'ready'
        AND ({" OR ".join(keyword_conditions)})
        {document_filter_sql}
        LIMIT %s
    """

    params = (
        params
        + document_filter_params
        + [limit]
    )

    cursor.execute(query, params)

    return cursor.fetchall()


def merge_retrieval_results(vector_rows, keyword_rows):
    merged = {}

    for row in vector_rows:
        chunk_id = row[0]

        merged[chunk_id] = {
            "chunk_id": row[0],
            "document_id": row[1],
            "filename": row[2],
            "chunk_number": row[3],
            "content": row[4],
            "vector_distance": float(row[5]),
            "from_vector": True,
            "from_keyword": False,
            "retrieval_type": "direct"
        }

    for row in keyword_rows:
        chunk_id = row[0]

        if chunk_id in merged:
            merged[chunk_id]["from_keyword"] = True
        else:
            merged[chunk_id] = {
                "chunk_id": row[0],
                "document_id": row[1],
                "filename": row[2],
                "chunk_number": row[3],
                "content": row[4],
                "vector_distance": float(row[5]),
                "from_vector": False,
                "from_keyword": True,
                "retrieval_type": "direct"
            }

    return list(merged.values())


def rerank_chunks(chunks, question):
    question_words = set(tokenize(question))
    reranked = []

    for chunk in chunks:
        content_words = set(tokenize(chunk["content"]))

        keyword_hits = len(question_words.intersection(content_words))
        vector_score = 1 / (1 + chunk["vector_distance"])

        source_bonus = 0

        if chunk["from_vector"]:
            source_bonus += 0.5

        if chunk["from_keyword"]:
            source_bonus += 1.0

        exact_phrase_bonus = 0

        if question.lower().strip() in chunk["content"].lower():
            exact_phrase_bonus = 3.0

        final_score = keyword_hits + vector_score + source_bonus + exact_phrase_bonus

        chunk["keyword_hits"] = keyword_hits
        chunk["vector_score"] = vector_score
        chunk["final_score"] = final_score

        reranked.append(chunk)

    reranked.sort(key=lambda item: item["final_score"], reverse=True)

    return reranked


def fetch_neighbor_chunks(cursor, ranked_chunks):
    neighbor_map = {}

    for chunk in ranked_chunks:
        document_id = chunk["document_id"]
        center_chunk_number = chunk["chunk_number"]

        start_chunk = max(1, center_chunk_number - NEIGHBOR_WINDOW)
        end_chunk = center_chunk_number + NEIGHBOR_WINDOW

        cursor.execute(
            """
            SELECT 
                c.chunk_id,
                c.document_id,
                d.original_filename,
                c.chunk_number,
                c.content
            FROM document_chunks c
            JOIN documents d
            ON c.document_id = d.document_id
            WHERE c.document_id = %s
              AND c.chunk_number BETWEEN %s AND %s
              AND d.indexing_status = 'ready'
            ORDER BY c.chunk_number ASC
            """,
            (
                document_id,
                start_chunk,
                end_chunk
            )
        )

        rows = cursor.fetchall()

        for row in rows:
            chunk_id = row[0]

            if chunk_id not in neighbor_map:
                neighbor_map[chunk_id] = {
                    "chunk_id": row[0],
                    "document_id": row[1],
                    "filename": row[2],
                    "chunk_number": row[3],
                    "content": row[4],
                    "vector_distance": chunk["vector_distance"],
                    "from_vector": chunk["from_vector"],
                    "from_keyword": chunk["from_keyword"],
                    "retrieval_type": "neighbor_expansion",
                    "keyword_hits": chunk.get("keyword_hits", 0),
                    "vector_score": chunk.get("vector_score", 0),
                    "final_score": chunk.get("final_score", 0)
                }

    return list(neighbor_map.values())


def build_context_chunks_for_document(cursor, ranked_chunks, dynamic_limit):
    direct_chunks = ranked_chunks[:PER_DOCUMENT_DIRECT_TOP_K]
    neighbor_chunks = fetch_neighbor_chunks(cursor, direct_chunks)

    combined = {}

    for chunk in direct_chunks + neighbor_chunks:
        combined[chunk["chunk_id"]] = chunk

    combined_chunks = list(combined.values())

    combined_chunks.sort(
        key=lambda item: item["chunk_number"]
    )

    return combined_chunks[:dynamic_limit]


def retrieve_relevant_chunks(question, document_ids=None):
    conn = get_db_connection()
    cursor = conn.cursor()

    requested_document_ids = document_ids or []

    if requested_document_ids:
        target_document_ids = requested_document_ids
    else:
        target_document_ids = get_ready_document_ids()

    if not target_document_ids:
        cursor.close()
        conn.close()
        return []

    per_document_limit = max(
        3,
        min(
            PER_DOCUMENT_CONTEXT_LIMIT,
            MAX_CONTEXT_CHUNKS_TOTAL // max(1, len(target_document_ids))
        )
    )

    final_chunks = []

    for document_id in target_document_ids:
        vector_rows = get_vector_rows(
            cursor,
            question,
            [document_id],
            PER_DOCUMENT_VECTOR_TOP_K
        )

        keyword_rows = get_keyword_rows(
            cursor,
            question,
            [document_id],
            PER_DOCUMENT_KEYWORD_TOP_K
        )

        merged_chunks = merge_retrieval_results(vector_rows, keyword_rows)
        ranked_chunks = rerank_chunks(merged_chunks, question)

        if not ranked_chunks:
            continue

        document_context_chunks = build_context_chunks_for_document(
            cursor,
            ranked_chunks,
            per_document_limit
        )

        final_chunks.extend(document_context_chunks)

    cursor.close()
    conn.close()

    return final_chunks[:MAX_CONTEXT_CHUNKS_TOTAL]


def create_session_if_needed(session_id, first_question):
    conn = get_db_connection()
    cursor = conn.cursor()

    if session_id is not None:
        cursor.close()
        conn.close()
        return session_id

    title = first_question[:60]

    cursor.execute(
        """
        INSERT INTO chat_sessions(title)
        VALUES (%s)
        RETURNING session_id
        """,
        (title,)
    )

    new_session_id = cursor.fetchone()[0]

    conn.commit()
    cursor.close()
    conn.close()

    return new_session_id


def save_message(session_id, role, message):
    conn = get_db_connection()
    cursor = conn.cursor()

    cursor.execute(
        """
        INSERT INTO chat_messages(session_id, role, message)
        VALUES (%s, %s, %s)
        """,
        (
            session_id,
            role,
            message
        )
    )

    conn.commit()
    cursor.close()
    conn.close()


def generate_answer(question, chunks):
    context_parts = []

    for index, chunk in enumerate(chunks, start=1):
        context_parts.append(
            f"Source {index}\n"
            f"Document ID: {chunk['document_id']}\n"
            f"File: {chunk['filename']}\n"
            f"Chunk: {chunk['chunk_number']}\n"
            f"Content:\n{chunk['content']}"
        )

    context = "\n\n-------------------------\n\n".join(context_parts)

    system_prompt = """
You are a strict enterprise knowledge-base assistant.

Use only the given knowledge base content.

Rules:
- Answer only from the given content.
- Do not use outside knowledge.
- Do not guess missing facts.
- Do not invent names, numbers, dates, policies, rules, or facts.
- If multiple documents are provided, consider all documents.
- Do not ignore a document unless no relevant content from that document is present.
- If different documents provide different information, mention it document-wise.
- If multiple related points are present, include all relevant points.
- If the exact answer is spread across multiple chunks, combine them into one complete answer.
- If the question is vague, summarize the most relevant information from all provided documents.
- Keep the answer clear, professional, and structured.
- Do not say "based on context" or "retrieved chunks".
"""

    user_prompt = f"""
Knowledge base content:
{context}

User question:
{question}

Give the final answer only.
"""

    response = requests.post(
        "http://localhost:11434/api/chat",
        json={
            "model": CHAT_MODEL,
            "messages": [
                {
                    "role": "system",
                    "content": system_prompt
                },
                {
                    "role": "user",
                    "content": user_prompt
                }
            ],
            "stream": False
        },
        timeout=180
    )

    if response.status_code != 200:
        print(response.text)
        raise Exception("Chat model failed")

    return response.json()["message"]["content"].strip()


@app.on_event("startup")
def startup_event():
    ensure_schema_updates()
    load_knowledge_folder_from_db()
    ensure_folders()

    print("\nKB Chatbot backend started.")
    print("Knowledge base folder:", os.path.abspath(KNOWLEDGE_BASE_FOLDER))
    print("Hourly indexing enabled. Manual indexing also available.")

    watcher_thread = threading.Thread(
        target=hourly_knowledge_base_watcher,
        daemon=True
    )

    watcher_thread.start()


@app.get("/")
def home():
    return FileResponse(UI_FILE)


@app.post("/login")
def login(request: LoginRequest, response: Response):
    conn = get_db_connection()
    cursor = conn.cursor()

    cursor.execute(
        """
        SELECT employee_id, name, email, role
        FROM employees
        WHERE email = %s
          AND password_hash = %s
        """,
        (
            request.email.strip().lower(),
            hash_password(request.password)
        )
    )

    row = cursor.fetchone()

    if not row:
        cursor.close()
        conn.close()
        raise HTTPException(status_code=401, detail="Invalid email or password")

    session_token = str(uuid.uuid4())

    cursor.execute(
        """
        INSERT INTO employee_sessions(employee_id, session_token)
        VALUES (%s, %s)
        """,
        (
            row[0],
            session_token
        )
    )

    conn.commit()
    cursor.close()
    conn.close()

    response.set_cookie(
        key="kb_session_token",
        value=session_token,
        httponly=True,
        samesite="lax"
    )

    return {
        "message": "Login successful",
        "employee": {
            "employee_id": row[0],
            "name": row[1],
            "email": row[2],
            "role": row[3]
        }
    }


@app.post("/logout")
def logout(request: Request, response: Response):
    token = get_session_token(request)

    if token:
        conn = get_db_connection()
        cursor = conn.cursor()

        cursor.execute(
            """
            DELETE FROM employee_sessions
            WHERE session_token = %s
            """,
            (token,)
        )

        conn.commit()
        cursor.close()
        conn.close()

    response.delete_cookie("kb_session_token")

    return {
        "message": "Logged out"
    }


@app.get("/me")
def me(request: Request):
    employee = require_login(request)

    return {
        "employee": employee
    }


@app.get("/health")
def health_check():
    return {
        "status": "Backend is running"
    }


@app.get("/knowledge-folder")
def get_knowledge_folder(request: Request):
    require_login(request)

    return {
        "folder_path": KNOWLEDGE_BASE_FOLDER
    }


@app.post("/set-knowledge-folder")
def set_knowledge_folder(request: FolderRequest, http_request: Request):
    require_login(http_request)

    global KNOWLEDGE_BASE_FOLDER
    global UPLOAD_FOLDER

    folder_path = request.folder_path.strip()

    if not os.path.isdir(folder_path):
        return {
            "success": False,
            "message": "Folder path does not exist.",
            "folder_path": folder_path
        }

    KNOWLEDGE_BASE_FOLDER = folder_path
    UPLOAD_FOLDER = KNOWLEDGE_BASE_FOLDER

    set_setting("knowledge_base_folder", folder_path)
    ensure_folders()

    return {
        "success": True,
        "message": "Knowledge folder saved successfully.",
        "folder_path": KNOWLEDGE_BASE_FOLDER
    }


@app.post("/index-now")
def index_now(request: Request, background_tasks: BackgroundTasks, force: bool = False):
    require_login(request)

    background_tasks.add_task(scan_knowledge_base_once, force)

    return {
        "message": "Manual indexing started in background.",
        "force": force
    }


@app.post("/upload")
async def upload_documents(request: Request, background_tasks: BackgroundTasks):
    require_login(request)
    ensure_folders()

    form = await request.form()
    uploaded_files = form.getlist("files")

    results = []

    if not uploaded_files:
        return {
            "message": "No files received.",
            "results": []
        }

    for file in uploaded_files:
        safe_filename = os.path.basename(file.filename)
        save_path = os.path.join(UPLOAD_FOLDER, safe_filename)

        with open(save_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)

        if not is_supported_file(save_path):
            results.append(
                {
                    "filename": safe_filename,
                    "status": "unsupported",
                    "message": "Unsupported file type.",
                    "scheduled": False
                }
            )
            continue

        result = queue_file_for_indexing(
            file_path=save_path,
            source_type="knowledge_base",
            force_reindex=False
        )

        results.append(result)

        if result.get("scheduled"):
            background_tasks.add_task(
                process_document_indexing,
                result["document_id"],
                os.path.abspath(save_path)
            )

    return {
        "message": "Upload completed.",
        "results": results
    }


@app.get("/documents")
def list_documents(request: Request):
    require_login(request)

    conn = get_db_connection()
    cursor = conn.cursor()

    cursor.execute(
        """
        SELECT 
            d.document_id,
            d.original_filename,
            d.source_type,
            d.version,
            d.indexing_status,
            d.error_message,
            CASE 
                WHEN d.chunk_count IS NULL OR d.chunk_count = 0
                THEN COUNT(c.chunk_id)
                ELSE d.chunk_count
            END AS chunks,
            d.indexed_at,
            d.updated_at
        FROM documents d
        LEFT JOIN document_chunks c
        ON d.document_id = c.document_id
        GROUP BY 
            d.document_id,
            d.original_filename,
            d.source_type,
            d.version,
            d.indexing_status,
            d.error_message,
            d.chunk_count,
            d.indexed_at,
            d.updated_at
        ORDER BY d.updated_at DESC;
        """
    )

    rows = cursor.fetchall()

    cursor.close()
    conn.close()

    return {
        "documents": [
            {
                "document_id": row[0],
                "filename": row[1],
                "source_type": row[2],
                "version": row[3],
                "indexing_status": row[4],
                "error_message": row[5],
                "chunks": row[6],
                "indexed_at": str(row[7]),
                "updated_at": str(row[8])
            }
            for row in rows
        ]
    }


@app.get("/download-document/{document_id}")
def download_document(document_id: int, request: Request):
    require_login(request)

    conn = get_db_connection()
    cursor = conn.cursor()

    cursor.execute(
        """
        SELECT original_filename, file_path
        FROM documents
        WHERE document_id = %s
        """,
        (document_id,)
    )

    row = cursor.fetchone()

    cursor.close()
    conn.close()

    if not row:
        raise HTTPException(status_code=404, detail="Document not found")

    filename, file_path = row

    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="File not found on server")

    return FileResponse(
        file_path,
        filename=filename,
        media_type="application/octet-stream"
    )


@app.post("/retrieve")
def retrieve(request_data: RetrieveRequest, request: Request):
    require_login(request)

    chunks = retrieve_relevant_chunks(
        request_data.question,
        request_data.document_ids
    )

    return {
        "question": request_data.question,
        "document_ids": request_data.document_ids,
        "search_scope": "selected_documents" if request_data.document_ids else "all_ready_documents",
        "chunks": [
            {
                "rank": index + 1,
                "document_id": chunk["document_id"],
                "filename": chunk["filename"],
                "chunk_number": chunk["chunk_number"],
                "retrieval_type": chunk["retrieval_type"],
                "final_score": round(chunk["final_score"], 4),
                "keyword_hits": chunk["keyword_hits"],
                "from_vector": chunk["from_vector"],
                "from_keyword": chunk["from_keyword"],
                "content_preview": chunk["content"][:700]
            }
            for index, chunk in enumerate(chunks)
        ]
    }


@app.post("/chat")
def chat(request_data: ChatRequest, request: Request):
    require_login(request)

    session_id = create_session_if_needed(
        request_data.session_id,
        request_data.question
    )

    save_message(session_id, "user", request_data.question)

    chunks = retrieve_relevant_chunks(
        request_data.question,
        request_data.document_ids
    )

    if not chunks:
        answer = "No ready and relevant knowledge base content was found for the selected document scope."
        save_message(session_id, "assistant", answer)

        return {
            "session_id": session_id,
            "answer": answer,
            "sources": []
        }

    answer = generate_answer(request_data.question, chunks)

    save_message(session_id, "assistant", answer)

    unique_sources = {}

    for chunk in chunks:
        unique_sources[chunk["document_id"]] = {
            "document_id": chunk["document_id"],
            "filename": chunk["filename"],
            "download_url": f"/download-document/{chunk['document_id']}"
        }

    return {
        "session_id": session_id,
        "answer": answer,
        "search_scope": "selected_documents" if request_data.document_ids else "all_ready_documents",
        "document_ids": request_data.document_ids,
        "sources": list(unique_sources.values())
    }


@app.get("/sessions")
def get_sessions(request: Request):
    require_login(request)

    conn = get_db_connection()
    cursor = conn.cursor()

    cursor.execute(
        """
        SELECT session_id, title, created_at
        FROM chat_sessions
        ORDER BY created_at DESC;
        """
    )

    rows = cursor.fetchall()

    cursor.close()
    conn.close()

    return {
        "sessions": [
            {
                "session_id": row[0],
                "title": row[1],
                "created_at": str(row[2])
            }
            for row in rows
        ]
    }


@app.get("/sessions/{session_id}/messages")
def get_session_messages(session_id: int, request: Request):
    require_login(request)

    conn = get_db_connection()
    cursor = conn.cursor()

    cursor.execute(
        """
        SELECT role, message, created_at
        FROM chat_messages
        WHERE session_id = %s
        ORDER BY created_at ASC;
        """,
        (session_id,)
    )

    rows = cursor.fetchall()

    cursor.close()
    conn.close()

    return {
        "messages": [
            {
                "role": row[0],
                "message": row[1],
                "created_at": str(row[2])
            }
            for row in rows
        ]
    } 