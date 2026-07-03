import os
import hashlib
import re
import csv
from datetime import datetime

import fitz
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from app.config.settings import (
    MAX_CHUNK_SIZE,
    CHUNK_OVERLAP
)


def calculate_file_hash(file_path):
    sha256 = hashlib.sha256()

    with open(file_path, "rb") as file:
        for block in iter(lambda: file.read(1024 * 1024), b""):
            sha256.update(block)

    return sha256.hexdigest()


def get_file_metadata(file_path):
    stat = os.stat(file_path)

    return {
        "file_size": stat.st_size,
        "last_modified": datetime.fromtimestamp(stat.st_mtime)
    }


def clean_text(text):
    if not text:
        return ""

    text = text.replace("\x00", "")
    text = text.replace("\r", "")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)

    return text.strip()


def tokenize(text):
    text = text.lower()
    text = re.sub(r"[^a-z0-9\s]", " ", text)

    stopwords = {
        "the", "is", "are", "was", "were", "a", "an", "of", "for", "to",
        "in", "on", "and", "or", "with", "from", "by", "as", "at", "this",
        "that", "it", "be", "can", "what", "which", "who", "when", "where",
        "how", "tell", "me", "about", "give", "list", "explain", "show",
        "please", "kindly", "all", "any"
    }

    return [word for word in text.split() if word not in stopwords]


def extract_pdf(file_path):
    document = fitz.open(file_path)
    text = ""

    for page_number, page in enumerate(document, start=1):
        text += f"\n\nPage {page_number}\n{page.get_text()}"

    document.close()
    return text


def extract_docx(file_path):
    doc = Document(file_path)
    text = ""

    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            text += paragraph.text.strip() + "\n"

    for table in doc.tables:
        for row in table.rows:
            row_values = []

            for cell in row.cells:
                if cell.text.strip():
                    row_values.append(cell.text.strip())

            if row_values:
                text += " | ".join(row_values) + "\n"

    return text


def extract_xlsx(file_path):
    workbook = load_workbook(file_path, data_only=True)
    text = ""

    for sheet in workbook.worksheets:
        text += f"\n\nSheet: {sheet.title}\n"

        for row in sheet.iter_rows(values_only=True):
            values = []

            for cell in row:
                if cell is not None:
                    values.append(str(cell))

            if values:
                text += " | ".join(values) + "\n"

    return text


def extract_csv(file_path):
    text = ""

    with open(file_path, "r", encoding="utf-8", errors="ignore") as file:
        reader = csv.reader(file)

        for row in reader:
            text += " | ".join(row) + "\n"

    return text


def extract_pptx(file_path):
    presentation = Presentation(file_path)
    text = ""

    for slide_number, slide in enumerate(presentation.slides, start=1):
        text += f"\n\nSlide {slide_number}\n"

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text += shape.text.strip() + "\n"

    return text


def extract_text_file(file_path):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as file:
        return file.read()


def extract_text(file_path):
    extension = file_path.lower().split(".")[-1]

    if extension == "pdf":
        return extract_pdf(file_path)

    if extension == "docx":
        return extract_docx(file_path)

    if extension == "xlsx":
        return extract_xlsx(file_path)

    if extension == "csv":
        return extract_csv(file_path)

    if extension == "pptx":
        return extract_pptx(file_path)

    if extension in ["txt", "md"]:
        return extract_text_file(file_path)

    raise ValueError("Unsupported file type: " + extension)


def is_supported_file(file_path):
    supported_extensions = [
        "pdf",
        "docx",
        "xlsx",
        "csv",
        "pptx",
        "txt",
        "md"
    ]

    extension = file_path.lower().split(".")[-1]

    return extension in supported_extensions


def recursive_split(text, separators):
    if len(text) <= MAX_CHUNK_SIZE:
        return [text]

    if not separators:
        return [
            text[i:i + MAX_CHUNK_SIZE]
            for i in range(0, len(text), MAX_CHUNK_SIZE)
        ]

    separator = separators[0]
    parts = text.split(separator)

    if len(parts) == 1:
        return recursive_split(text, separators[1:])

    chunks = []
    current = ""

    for part in parts:
        part = part.strip()

        if not part:
            continue

        candidate = current + separator + part if current else part

        if len(candidate) <= MAX_CHUNK_SIZE:
            current = candidate
        else:
            if current:
                chunks.extend(recursive_split(current, separators[1:]))

            current = part

    if current:
        chunks.extend(recursive_split(current, separators[1:]))

    return chunks


def add_overlap(chunks):
    final_chunks = []

    for chunk in chunks:
        chunk = chunk.strip()

        if not chunk:
            continue

        if not final_chunks:
            final_chunks.append(chunk)
        else:
            previous = final_chunks[-1]
            overlap_text = previous[-CHUNK_OVERLAP:]
            combined = overlap_text + "\n" + chunk

            if len(combined) <= MAX_CHUNK_SIZE + CHUNK_OVERLAP:
                final_chunks.append(combined.strip())
            else:
                final_chunks.append(chunk)

    return final_chunks


def split_text_into_chunks(text):
    text = clean_text(text)

    separators = [
        "\n\n",
        "\n",
        ". ",
        "; ",
        ", ",
        " "
    ]

    raw_chunks = recursive_split(text, separators)

    safe_chunks = []

    for chunk in raw_chunks:
        if len(chunk) <= MAX_CHUNK_SIZE:
            safe_chunks.append(chunk)
        else:
            for i in range(0, len(chunk), MAX_CHUNK_SIZE):
                safe_chunks.append(chunk[i:i + MAX_CHUNK_SIZE])

    final_chunks = add_overlap(safe_chunks)

    return [
        chunk.strip()
        for chunk in final_chunks
        if len(chunk.strip()) > 30
    ]